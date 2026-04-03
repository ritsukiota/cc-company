#!/usr/bin/env python3
"""
PowerPoint互換 PPTXジェネレーター (python-pptx)
Usage: python3 generate-pptx-py.py <input.json> <output.pptx>
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── レイアウト定数 ───
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN_TOP = Inches(0.4)
MARGIN_BOTTOM = Inches(0.3)
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_H = SLIDE_H - MARGIN_TOP - MARGIN_BOTTOM
TITLE_H = Inches(0.65)
BODY_Y = MARGIN_TOP + TITLE_H + Inches(0.1)
BODY_H = CONTENT_H - TITLE_H - Inches(0.1)

# ─── カラースキーム ───
COLOR_SCHEMES = {
    "navy-professional": {
        "primary": "1B2A4A", "secondary": "2D4A7A",
        "accent1": "3B82F6", "accent2": "10B981", "accent3": "F59E0B",
        "textDark": "1F2937", "textLight": "FFFFFF", "textMuted": "6B7280",
        "background": "FFFFFF", "backgroundAlt": "F3F4F6",
        "positive": "10B981", "negative": "EF4444", "warning": "F59E0B", "border": "E5E7EB",
    },
    "dark-executive": {
        "primary": "1F2937", "secondary": "374151",
        "accent1": "60A5FA", "accent2": "34D399", "accent3": "FBBF24",
        "textDark": "F9FAFB", "textLight": "FFFFFF", "textMuted": "9CA3AF",
        "background": "111827", "backgroundAlt": "1F2937",
        "positive": "34D399", "negative": "F87171", "warning": "FBBF24", "border": "374151",
    },
    "clean-minimal": {
        "primary": "374151", "secondary": "6B7280",
        "accent1": "2563EB", "accent2": "059669", "accent3": "D97706",
        "textDark": "111827", "textLight": "FFFFFF", "textMuted": "9CA3AF",
        "background": "FFFFFF", "backgroundAlt": "F9FAFB",
        "positive": "059669", "negative": "DC2626", "warning": "D97706", "border": "E5E7EB",
    },
    "warm-corporate": {
        "primary": "991B1B", "secondary": "B91C1C",
        "accent1": "DC2626", "accent2": "059669", "accent3": "D97706",
        "textDark": "1F2937", "textLight": "FFFFFF", "textMuted": "6B7280",
        "background": "FFFBEB", "backgroundAlt": "FEF3C7",
        "positive": "059669", "negative": "DC2626", "warning": "D97706", "border": "E5E7EB",
    },
}

FONT = "Hiragino Kaku Gothic ProN"


def rgb(hex_str):
    return RGBColor.from_string(hex_str)


def add_bg_rect(slide, color):
    """全面背景矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color="1F2937", align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name=FONT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
        try:
            from pptx.oxml.ns import qn
            txBody = txBox.text_frame._txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            anchor_map = {
                MSO_ANCHOR.TOP: 't',
                MSO_ANCHOR.MIDDLE: 'ctr',
                MSO_ANCHOR.BOTTOM: 'b',
            }
            bodyPr.set('anchor', anchor_map.get(anchor, 't'))
        except:
            pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_slide_title(slide, title, colors):
    """スライドタイトル + アクセントライン"""
    add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, TITLE_H,
                title, font_size=28, bold=True, color=colors["primary"],
                anchor=MSO_ANCHOR.BOTTOM)
    # アクセントライン
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_LEFT, MARGIN_TOP + TITLE_H + Emu(18288),
        CONTENT_W, Emu(27432))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colors["accent1"])
    shape.line.fill.background()


def add_page_number(slide, num, colors):
    add_textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.4),
                Inches(0.5), Inches(0.3), str(num),
                font_size=10, color=colors["textMuted"], align=PP_ALIGN.RIGHT)


def add_rect(slide, left, top, width, height, fill_color, radius=False):
    """矩形シェイプ"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.fill.background()
    return shape


def add_bullets(slide, left, top, width, height, items, font_size=14, color="1F2937"):
    """箇条書きテキスト"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(color)
        p.font.name = FONT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        # 箇条書きマーカー
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        pPr.append(buChar)
        pPr.set('indent', str(Pt(14)))
        pPr.set('marL', str(Pt(18)))
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths_ratio, colors, row_height=Inches(0.45)):
    """テーブルを追加"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    if n_rows == 0 or n_cols == 0:
        return

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table

    # カラム幅設定
    total_ratio = sum(col_widths_ratio)
    for ci, ratio in enumerate(col_widths_ratio):
        table.columns[ci].width = int(width * ratio / total_ratio)

    for ri, row in enumerate(rows_data):
        table.rows[ri].height = row_height
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text

            # セルの書式
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12 if ri == 0 else 11)
            p.font.bold = (ri == 0)
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

            if ri == 0:
                p.font.color.rgb = rgb(colors["textLight"])
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["primary"])
            else:
                p.font.color.rgb = rgb(colors["textDark"])
                cell.fill.solid()
                fill_c = colors["backgroundAlt"] if ri % 2 == 0 else colors["background"]
                cell.fill.fore_color.rgb = rgb(fill_c)

    return table_shape


# ─── スライドレンダラー ───

def render_title(prs, data, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_bg_rect(slide, colors["primary"])
    add_textbox(slide, MARGIN_LEFT, Inches(1.8), CONTENT_W, Inches(1.5),
                data["title"], font_size=36, bold=True, color=colors["textLight"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if data.get("subtitle"):
        add_textbox(slide, MARGIN_LEFT, Inches(3.4), CONTENT_W, Inches(0.8),
                    data["subtitle"], font_size=18, color=colors["textLight"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if data.get("authorLine"):
        add_textbox(slide, MARGIN_LEFT, Inches(5.5), CONTENT_W, Inches(0.5),
                    data["authorLine"], font_size=14, color=colors["textMuted"],
                    align=PP_ALIGN.CENTER)
    # 下部アクセントライン
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), colors["accent1"])


def render_executive_summary(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    has_metric = "keyMetric" in data and data["keyMetric"]
    bullet_w = int(CONTENT_W * 0.62) if has_metric else CONTENT_W

    add_bullets(slide, MARGIN_LEFT, BODY_Y, bullet_w, BODY_H,
                data["bullets"], font_size=14, color=colors["textDark"])

    if has_metric:
        mx = MARGIN_LEFT + int(CONTENT_W * 0.65)
        mw = int(CONTENT_W * 0.35)
        add_rect(slide, mx, BODY_Y + Inches(0.3), mw, Inches(2.2),
                 colors["backgroundAlt"], radius=True)
        add_textbox(slide, mx, BODY_Y + Inches(0.6), mw, Inches(1.0),
                    data["keyMetric"]["value"], font_size=48, bold=True,
                    color=colors["accent1"], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, mx, BODY_Y + Inches(1.6), mw, Inches(0.5),
                    data["keyMetric"]["label"], font_size=14,
                    color=colors["textMuted"], align=PP_ALIGN.CENTER)


def render_agenda(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    items = data["items"]
    item_h = min(Inches(0.9), int(BODY_H / len(items)))

    for i, item in enumerate(items):
        y = BODY_Y + i * item_h
        # 番号
        add_textbox(slide, MARGIN_LEFT, y, Inches(0.8), item_h,
                    f"{item['number']:02d}", font_size=28, bold=True,
                    color=colors["accent1"], anchor=MSO_ANCHOR.MIDDLE)
        # テキスト
        text_h = int(item_h * 0.55) if item.get("description") else item_h
        add_textbox(slide, MARGIN_LEFT + Inches(0.9), y, CONTENT_W - Inches(1.0), text_h,
                    item["text"], font_size=18, bold=True, color=colors["textDark"],
                    anchor=MSO_ANCHOR.BOTTOM if item.get("description") else MSO_ANCHOR.MIDDLE)
        if item.get("description"):
            add_textbox(slide, MARGIN_LEFT + Inches(0.9), y + int(item_h * 0.5),
                        CONTENT_W - Inches(1.0), int(item_h * 0.45),
                        item["description"], font_size=12, color=colors["textMuted"],
                        anchor=MSO_ANCHOR.TOP)
        # 区切り線
        if i < len(items) - 1:
            add_rect(slide, MARGIN_LEFT + Inches(0.9), y + item_h - Emu(18288),
                     CONTENT_W - Inches(1.0), Emu(9144), colors["border"])


def render_market_analysis(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    table_y = BODY_Y
    if data.get("insight"):
        add_rect(slide, MARGIN_LEFT, BODY_Y, CONTENT_W, Inches(0.5),
                 colors["accent1"], radius=True)
        add_textbox(slide, MARGIN_LEFT + Inches(0.15), BODY_Y,
                    CONTENT_W - Inches(0.3), Inches(0.5),
                    data["insight"], font_size=12, bold=True,
                    color=colors["textLight"], anchor=MSO_ANCHOR.MIDDLE)
        table_y = BODY_Y + Inches(0.65)

    if data.get("data"):
        headers = data["data"]["headers"]
        rows = [headers] + data["data"]["rows"]
        n_cols = len(headers)
        col_ratios = [1.0] * n_cols
        add_table(slide, MARGIN_LEFT, table_y, CONTENT_W, rows, col_ratios, colors)

    if data.get("source"):
        add_textbox(slide, MARGIN_LEFT, SLIDE_H - Inches(0.7), CONTENT_W, Inches(0.3),
                    data["source"], font_size=10, color=colors["textMuted"])


def render_swot(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    gap = Inches(0.12)
    so_what_h = Inches(0.55) if data.get("soWhat") else 0
    quad_w = (CONTENT_W - gap) // 2
    quad_h = (BODY_H - gap - so_what_h) // 2
    label_h = Inches(0.38)

    quadrants = [
        ("Strengths", data.get("strengths", []), MARGIN_LEFT, BODY_Y, colors["accent1"]),
        ("Weaknesses", data.get("weaknesses", []), MARGIN_LEFT + quad_w + gap, BODY_Y, colors["textMuted"]),
        ("Opportunities", data.get("opportunities", []), MARGIN_LEFT, BODY_Y + quad_h + gap, colors["accent2"]),
        ("Threats", data.get("threats", []), MARGIN_LEFT + quad_w + gap, BODY_Y + quad_h + gap, colors["warning"]),
    ]

    for label, items, x, y, color in quadrants:
        add_rect(slide, x, y, quad_w, quad_h, colors["backgroundAlt"], radius=True)
        add_rect(slide, x, y, quad_w, label_h, color)
        add_textbox(slide, x + Inches(0.1), y, quad_w - Inches(0.2), label_h,
                    label, font_size=13, bold=True, color=colors["textLight"],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(slide, x + Inches(0.15), y + label_h + Inches(0.05),
                    quad_w - Inches(0.3), quad_h - label_h - Inches(0.1),
                    items, font_size=11, color=colors["textDark"])

    if data.get("soWhat"):
        so_y = BODY_Y + quad_h * 2 + gap + Inches(0.08)
        add_rect(slide, MARGIN_LEFT, so_y, CONTENT_W, Inches(0.45), colors["accent1"], radius=True)
        add_textbox(slide, MARGIN_LEFT + Inches(0.15), so_y,
                    CONTENT_W - Inches(0.3), Inches(0.45),
                    f"So What?  {data['soWhat']}", font_size=12, bold=True,
                    color=colors["textLight"], anchor=MSO_ANCHOR.MIDDLE)


def render_competitive_landscape(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    headers = [""] + data["criteria"]
    rows_data = [[comp] + scores for comp, scores in zip(data["competitors"], data["scores"])]
    all_rows = [headers] + rows_data
    n_cols = len(headers)
    col_ratios = [2.0] + [1.0] * (n_cols - 1)

    add_table(slide, MARGIN_LEFT, BODY_Y, CONTENT_W, all_rows, col_ratios, colors)

    if data.get("positioning"):
        add_textbox(slide, MARGIN_LEFT, SLIDE_H - Inches(0.9), CONTENT_W, Inches(0.4),
                    data["positioning"], font_size=12, color=colors["textMuted"])


def render_strategic_framework(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    fw_type = data.get("frameworkType", "process")
    if fw_type == "process":
        render_process(slide, data, colors)
    elif fw_type == "matrix":
        render_matrix(slide, data, colors)
    elif fw_type == "pyramid":
        render_pyramid(slide, data, colors)


def render_process(slide, data, colors):
    steps = data.get("steps", [])
    n = len(steps)
    gap = Inches(0.1)
    total_gap = int(gap * (n - 1))
    step_w = (CONTENT_W - total_gap) // n
    step_h = Inches(1.2)
    step_y = BODY_Y + Inches(0.5)
    step_colors = [colors["accent1"], colors["accent2"], colors["accent3"],
                   colors["primary"], colors["secondary"]]

    for i, step in enumerate(steps):
        x = MARGIN_LEFT + i * (step_w + gap)
        c = step_colors[i % len(step_colors)]
        add_rect(slide, x, step_y, step_w, step_h, c, radius=True)
        add_textbox(slide, x, step_y, step_w, step_h,
                    step["label"], font_size=13, bold=True, color=colors["textLight"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if step.get("description"):
            add_textbox(slide, x, step_y + step_h + Inches(0.15), step_w, Inches(1.5),
                        step["description"], font_size=10, color=colors["textDark"])
        if i < n - 1:
            add_textbox(slide, x + step_w - Emu(45720), step_y, gap + Emu(91440), step_h,
                        "\u25B6", font_size=18, color=colors["textMuted"],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_matrix(slide, data, colors):
    cx = MARGIN_LEFT + CONTENT_W // 2
    cy = BODY_Y + BODY_H // 2
    quad_w = int(CONTENT_W * 0.42)
    quad_h = int(BODY_H * 0.40)
    gap = Inches(0.15)
    quad_colors = [colors["accent1"], colors["accent2"], colors["accent3"], colors["textMuted"]]

    quads = [
        ("topLeft", cx - quad_w - gap // 2, cy - quad_h - gap // 2),
        ("topRight", cx + gap // 2, cy - quad_h - gap // 2),
        ("bottomLeft", cx - quad_w - gap // 2, cy + gap // 2),
        ("bottomRight", cx + gap // 2, cy + gap // 2),
    ]

    for idx, (key, x, y) in enumerate(quads):
        add_rect(slide, x, y, quad_w, quad_h, colors["backgroundAlt"], radius=True)
        if data.get("quadrants") and data["quadrants"].get(key):
            add_textbox(slide, x + Inches(0.15), y + Inches(0.15),
                        quad_w - Inches(0.3), quad_h - Inches(0.3),
                        data["quadrants"][key], font_size=12, color=colors["textDark"],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_pyramid(slide, data, colors):
    levels = data.get("levels", [])
    n = len(levels)
    pyramid_w = int(CONTENT_W * 0.5)
    pyramid_h = int(BODY_H * 0.85)
    level_h = pyramid_h // n
    start_x = MARGIN_LEFT + int(CONTENT_W * 0.05)
    start_y = BODY_Y + Inches(0.1)
    level_colors = [colors["accent1"], colors["accent2"], colors["accent3"]]

    for i, level in enumerate(levels):
        ratio = (i + 0.5) / n
        w = int(pyramid_w * (0.3 + ratio * 0.7))
        x = start_x + (pyramid_w - w) // 2
        y = start_y + (n - 1 - i) * level_h
        c = level_colors[i % len(level_colors)]
        add_rect(slide, x, y, w, level_h - Emu(73152), c, radius=True)
        add_textbox(slide, x, y, w, level_h - Emu(73152),
                    level["label"], font_size=13, bold=True, color=colors["textLight"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if level.get("description"):
            add_textbox(slide, start_x + pyramid_w + Inches(0.3), y,
                        CONTENT_W - pyramid_w - Inches(0.6), level_h - Emu(73152),
                        level["description"], font_size=11, color=colors["textDark"],
                        anchor=MSO_ANCHOR.MIDDLE)


def render_roadmap(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    phases = data.get("phases", [])
    n = len(phases)
    gap = Inches(0.08)
    phase_w = (CONTENT_W - int(gap * (n - 1))) // n
    chevron_h = Inches(1.0)
    chevron_y = BODY_Y + Inches(0.4)
    phase_colors = [colors["accent1"], colors["accent2"], colors["accent3"],
                    colors["primary"], colors["secondary"]]

    for i, phase in enumerate(phases):
        x = MARGIN_LEFT + i * (phase_w + gap)
        c = phase_colors[i % len(phase_colors)]

        # 期間ラベル
        add_textbox(slide, x, BODY_Y, phase_w, Inches(0.35),
                    phase["period"], font_size=10, color=colors["textMuted"],
                    align=PP_ALIGN.CENTER)
        # シェブロン
        add_rect(slide, x, chevron_y, phase_w, chevron_h, c, radius=True)
        add_textbox(slide, x, chevron_y, phase_w, chevron_h,
                    phase["name"], font_size=14, bold=True, color=colors["textLight"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # マイルストーン
        if phase.get("milestones"):
            add_bullets(slide, x, chevron_y + chevron_h + Inches(0.15),
                        phase_w, BODY_H - chevron_h - Inches(0.9),
                        phase["milestones"], font_size=10, color=colors["textDark"])
        # 矢印
        if i < n - 1:
            add_textbox(slide, x + phase_w - Emu(45720), chevron_y,
                        gap + Emu(91440), chevron_h,
                        "\u25B6", font_size=16, color=colors["textMuted"],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_recommendations(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    items = data.get("items", [])
    item_h = min(Inches(1.0), int(BODY_H / len(items)))
    priority_colors = {"high": colors["negative"], "medium": colors["warning"], "low": colors["accent1"]}
    priority_labels = {"high": "HIGH", "medium": "MED", "low": "LOW"}

    for i, item in enumerate(items):
        y = BODY_Y + i * item_h

        # 番号（丸）
        add_rect(slide, MARGIN_LEFT, y + Inches(0.1), Inches(0.5), Inches(0.5),
                 colors["accent1"], radius=True)
        add_textbox(slide, MARGIN_LEFT, y + Inches(0.1), Inches(0.5), Inches(0.5),
                    str(item["number"]), font_size=18, bold=True,
                    color=colors["textLight"], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

        # テキスト
        add_textbox(slide, MARGIN_LEFT + Inches(0.7), y + Inches(0.05),
                    CONTENT_W - Inches(2.5), int(item_h * 0.65),
                    item["text"], font_size=14, color=colors["textDark"],
                    anchor=MSO_ANCHOR.MIDDLE)

        # 優先度バッジ
        pc = priority_colors.get(item.get("priority"), colors["textMuted"])
        add_rect(slide, MARGIN_LEFT + CONTENT_W - Inches(1.5), y + Inches(0.15),
                 Inches(0.8), Inches(0.35), pc, radius=True)
        add_textbox(slide, MARGIN_LEFT + CONTENT_W - Inches(1.5), y + Inches(0.15),
                    Inches(0.8), Inches(0.35),
                    priority_labels.get(item.get("priority", ""), ""),
                    font_size=9, bold=True, color=colors["textLight"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 区切り線
        if i < len(items) - 1:
            add_rect(slide, MARGIN_LEFT + Inches(0.7), y + item_h - Emu(18288),
                     CONTENT_W - Inches(0.7), Emu(9144), colors["border"])


def render_next_steps(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    items = data.get("items", [])
    headers = ["\u30A2\u30AF\u30B7\u30E7\u30F3", "\u62C5\u5F53", "\u671F\u9650"]
    rows = [[item["action"], item["owner"], item["deadline"]] for item in items]
    all_rows = [headers] + rows
    col_ratios = [0.55, 0.22, 0.23]

    add_table(slide, MARGIN_LEFT, BODY_Y, CONTENT_W, all_rows, col_ratios, colors,
              row_height=Inches(0.5))


def render_financial_impact(prs, data, colors, pn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, data["title"], colors)
    add_page_number(slide, pn, colors)

    metrics = data.get("metrics", [])
    n = len(metrics)
    gap = Inches(0.3)
    card_w = (CONTENT_W - int(gap * (n - 1))) // n
    card_h = Inches(2.2)

    for i, metric in enumerate(metrics):
        x = MARGIN_LEFT + i * (card_w + gap)
        add_rect(slide, x, BODY_Y, card_w, card_h, colors["backgroundAlt"], radius=True)
        add_textbox(slide, x, BODY_Y + Inches(0.2), card_w, Inches(1.0),
                    metric["value"], font_size=48, bold=True, color=colors["accent1"],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, BODY_Y + Inches(1.2), card_w, Inches(0.4),
                    metric["label"], font_size=14, color=colors["textMuted"],
                    align=PP_ALIGN.CENTER)
        if metric.get("delta"):
            is_positive = "+" in metric["delta"] or "\u2191" in metric["delta"]
            add_textbox(slide, x, BODY_Y + Inches(1.6), card_w, Inches(0.4),
                        metric["delta"], font_size=14, bold=True,
                        color=colors["positive"] if is_positive else colors["negative"],
                        align=PP_ALIGN.CENTER)


def render_section_divider(prs, data, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, colors["primary"])
    add_textbox(slide, MARGIN_LEFT, Inches(2.2), CONTENT_W, Inches(1.5),
                data["title"], font_size=32, bold=True, color=colors["textLight"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if data.get("subtitle"):
        add_textbox(slide, MARGIN_LEFT, Inches(3.7), CONTENT_W, Inches(0.8),
                    data["subtitle"], font_size=16, color=colors["textMuted"],
                    align=PP_ALIGN.CENTER)
    add_rect(slide, int(SLIDE_W * 0.35), Inches(4.5), int(SLIDE_W * 0.3), Emu(36576),
             colors["accent1"])


# ─── ディスパッチャー ───

RENDERERS = {
    "title": lambda prs, d, c, _: render_title(prs, d, c),
    "executive-summary": render_executive_summary,
    "agenda": render_agenda,
    "market-analysis": render_market_analysis,
    "swot": render_swot,
    "competitive-landscape": render_competitive_landscape,
    "strategic-framework": render_strategic_framework,
    "roadmap": render_roadmap,
    "financial-impact": render_financial_impact,
    "recommendations": render_recommendations,
    "next-steps": render_next_steps,
    "section-divider": lambda prs, d, c, _: render_section_divider(prs, d, c),
}


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 generate-pptx-py.py <input.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    input_path, output_path = sys.argv[1], sys.argv[2]

    with open(input_path, "r", encoding="utf-8") as f:
        definition = json.load(f)

    meta = definition["meta"]
    slides = definition["slides"]
    scheme_name = meta.get("colorScheme", "navy-professional")
    colors = COLOR_SCHEMES.get(scheme_name, COLOR_SCHEMES["navy-professional"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(slides):
        renderer = RENDERERS.get(slide_data["type"])
        if renderer:
            renderer(prs, slide_data, colors, i + 1)
        else:
            print(f"Unknown slide type: {slide_data['type']}, skipping.", file=sys.stderr)

    prs.save(output_path)
    print(f"PPTX generated: {output_path}")


if __name__ == "__main__":
    main()
